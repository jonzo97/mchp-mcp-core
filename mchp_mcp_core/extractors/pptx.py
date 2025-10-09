"""
PPTX (PowerPoint) extraction.

Extracts text from PowerPoint presentations at slide level.
"""

import hashlib
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

from mchp_mcp_core.models.common import ExtractedChunk
from mchp_mcp_core.utils.logger import get_logger

logger = get_logger(__name__)


class PPTXExtractor:
    """
    Handles PPTX extraction at slide level.

    Features:
    - Slide-by-slide text extraction
    - Shape text aggregation
    - Metadata preservation

    Example:
        >>> extractor = PPTXExtractor()
        >>> chunks = extractor.extract_document("presentation.pptx", "doc_123")
    """

    def extract_document(self, pptx_path: str, document_id: str) -> List[ExtractedChunk]:
        """
        Extract text chunks from PPTX file (one chunk per slide).

        Args:
            pptx_path: Path to PPTX file
            document_id: Unique identifier for the document

        Returns:
            List of extracted chunks (one per slide)
        """
        logger.info(f"Extracting PPTX: {pptx_path}")
        chunks = []

        try:
            prs = Presentation(pptx_path)

            for i, slide in enumerate(prs.slides, start=1):
                # Extract text from all shapes
                text = "\n".join([
                    shape.text
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                ])
                text = text.strip()

                if not text:
                    logger.debug(f"Slide {i} has no text content, skipping")
                    continue

                metadata = {
                    "slide": i,
                    "total_slides": len(prs.slides),
                    "extraction_method": "python-pptx",
                    "filename": Path(pptx_path).name
                }

                chunk = ExtractedChunk(
                    chunk_id=self._generate_chunk_id(text, i),
                    content=text,
                    page_start=i,
                    page_end=i,
                    chunk_type="slide",
                    section_hierarchy=f"Slide {i}",
                    metadata=metadata
                )

                chunks.append(chunk)

            logger.info(f"Extracted {len(chunks)} slides from {Path(pptx_path).name}")

        except Exception as e:
            logger.error(f"Error processing PPTX {pptx_path}: {e}")
            raise

        return chunks

    def get_document_metadata(self, pptx_path: str) -> Dict[str, Any]:
        """
        Extract document metadata from PPTX.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            Dictionary with metadata
        """
        try:
            prs = Presentation(pptx_path)
            core_props = prs.core_properties

            metadata = {
                "total_slides": len(prs.slides),
                "title": core_props.title or "",
                "author": core_props.author or "",
                "subject": core_props.subject or "",
                "created": core_props.created.isoformat() if core_props.created else None,
                "modified": core_props.modified.isoformat() if core_props.modified else None,
                "filename": Path(pptx_path).name
            }

            return metadata

        except Exception as e:
            logger.error(f"Error extracting metadata from {pptx_path}: {e}")
            return {
                "total_slides": 0,
                "filename": Path(pptx_path).name,
                "error": str(e)
            }

    def _generate_chunk_id(self, content: str, slide_num: int) -> str:
        """Generate a unique chunk ID for a slide."""
        hash_input = f"{content[:100]}{slide_num}".encode()
        hash_val = hashlib.md5(hash_input).hexdigest()[:12]
        return f"slide_{slide_num}_{hash_val}"


__all__ = ["PPTXExtractor"]
